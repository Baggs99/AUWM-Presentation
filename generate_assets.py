"""Generate the PowerPoint deck and Word script for the presentation.

Run:
    python generate_assets.py

Produces:
    I_Know_Who_I_Am_Presentation.pptx
    I_Know_Who_I_Am_Full_Script.docx

Requires: python-pptx, python-docx
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

from docx import Document
from docx.shared import Pt as DocxPt, Inches as DocxInches, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH

# ---------------------------------------------------------------------------
# Shared content
# ---------------------------------------------------------------------------

TITLE = "I Know Who I Am"
SUBTITLE = "A practical theology of African urban worship music"
PRESENTER = "Daniel Baglini"
COURSE = "African Urban Worship Music"

THESIS = (
    "\u201cI Know Who I Am\u201d is not just a song about identity. It is a practical "
    "theological act that forms believers to live differently in the world. By teaching "
    "worshippers who they are before God, the song gives them confidence, dignity, and a "
    "sense of responsibility that can move belief into action."
)

SLIDES = [
    {
        "kind": "title",
        "eyebrow": "Yale Divinity School  \u00b7  Final presentation",
        "title": TITLE,
        "subtitle": SUBTITLE,
        "meta": [
            ("Song", "\u201cI Know Who I Am\u201d"),
            ("Artist", "Sinach"),
            ("Course", COURSE),
            ("Presenter", PRESENTER),
        ],
        "notes": (
            "Open simply. State the song, the artist, and the central claim. Tell the class "
            "you want them to leave with one idea: this song is not only a song about identity. "
            "It is a practical theological act that forms believers to live differently."
        ),
    },
    {
        "kind": "content",
        "eyebrow": "Introduction",
        "title": "More than a song. A practical theological act.",
        "bullets": [
            "Sung across Africa and in the diaspora, not just performed.",
            "Repetition and participation are formation, not decoration.",
            "The song trains worshippers to carry a new identity into daily life.",
        ],
        "notes": (
            "Anchor the audience. In rooms in Lagos, Nairobi, Accra, Houston, or London, "
            "people sing, repeat, and move. That repetition is practical theology in real time. "
            "It is theology as practice, training people to carry a new identity out into "
            "Monday morning."
        ),
    },
    {
        "kind": "content",
        "eyebrow": "Song and artist context",
        "title": "Sinach and a song that traveled.",
        "bullets": [
            "Sinach: Nigerian gospel artist, associated with Christ Embassy.",
            "Circulates in churches across Africa, the diaspora, and beyond.",
            "An example of African worship music moving transnationally.",
            "Built for participation, not performance.",
        ],
        "notes": (
            "Keep context tight. Identify Sinach as a Nigerian gospel artist. Note the song\u2019s "
            "spread into Kenyan, Ghanaian, South African, US, and UK churches. Emphasize that "
            "the song was designed to be sung by a whole room together, which matters for the "
            "practical theology argument that follows."
        ),
    },
    {
        "kind": "content",
        "eyebrow": "What the song does",
        "title": "Music as formation.",
        "bullets": [
            "Repetition: the chorus presses the same claims into memory.",
            "Call and response: leader sets the line, the room answers.",
            "Rhythm and groove: the song is moved to, not only listened to.",
            "Embodied worship: sung with the whole body, not only the head.",
        ],
        "notes": (
            "Walk through the four musical features. Make the turn from description to "
            "argument: these are not decorative features, they are the machinery that makes the "
            "song formative rather than merely expressive. Keep this slide tight so you can move "
            "into what the song is actually claiming."
        ),
    },
    {
        "kind": "lyrics",
        "eyebrow": "What the song says",
        "title": "Identity claimed in the first person.",
        "lyrics": [
            "\u201cI know who God says I am\u201d",
            "\u201cI am holy\u201d",
            "\u201cI am righteous\u201d",
            "\u201cI am so rich\u201d",
        ],
        "argument": (
            "These are first person affirmations made out loud and made together. The singer is "
            "not quoting someone else. They are claiming the identity themselves, and the room "
            "is claiming it with them. The identity being sung is public, communal, and on its "
            "way to becoming practice."
        ),
        "notes": (
            "Pause on this slide. Read the lines slowly. Point out that these are first person "
            "claims made publicly and communally. This is the pivot: from what the song says to "
            "what the song is doing as practical theology."
        ),
    },
    {
        "kind": "practical_theology",
        "eyebrow": "Practical theology",
        "title": "From worship to action.",
        "subtitle": (
            "What we are watching is practical theology: worship that forms believers to live "
            "differently in the world."
        ),
        "steps": [
            {
                "num": "01",
                "label": "Names identity",
                "body": (
                    "The song gives worshippers specific language about who they are before "
                    "God: holy, righteous, loved, chosen."
                ),
            },
            {
                "num": "02",
                "label": "Forms disposition",
                "body": (
                    "Weekly repetition shapes how people carry themselves. Confidence, dignity, "
                    "and a sense of what they owe each other become habits, not moods."
                ),
            },
            {
                "num": "03",
                "label": "Shapes practice",
                "body": (
                    "Identity that is sung becomes identity that is lived. Worship moves into "
                    "work, family, and public witness the rest of the week."
                ),
            },
        ],
        "notes": (
            "This is the center of the talk. Slow down here. Make the three moves explicit: "
            "name, form, shape. Emphasize that the logic is moving from what the song is saying "
            "in the room to how it shapes lives outside the room. This is worship as formation, "
            "not decoration."
        ),
    },
    {
        "kind": "content",
        "eyebrow": "Why it matters",
        "title": "Belief moves into action.",
        "bullets": [
            "Music is not the warm up. It can be the teaching.",
            "Moral agency: words sung in public raise the floor for how we live.",
            "Communal encouragement across geography, a shared confidence to carry.",
            "Public Christian witness: Sunday formation shows up on Monday.",
        ],
        "notes": (
            "Zoom out. Argue that in African urban worship, music is often doing the teaching. "
            "Make the belief in action point strongly: if I have stood in a room and said I am "
            "holy, I am loved, I have a harder time walking out and living as if none of that "
            "is true. Close by tying this to public Christian witness."
        ),
    },
    {
        "kind": "quote",
        "eyebrow": "Conclusion",
        "title": "Faith formed, faith lived.",
        "quote": (
            "\u201cI Know Who I Am\u201d is practical theology set to a beat. It does not only "
            "express faith. It trains people to live it, and it sends them out to live it together."
        ),
        "notes": (
            "Land the close. Restate the thesis in plain language: this song forms believers to "
            "carry a new identity, and a new responsibility, into their daily lives. Thank the "
            "class and invite questions."
        ),
    },
    {
        "kind": "bibliography",
        "eyebrow": "Bibliography",
        "title": "Sources",
        "entries": [
            "Ajose, Toyin. \u201cMultitracking the Spirit: Musicality, Prayer Chants and "
            "Pentecostal Spirituality in Nigeria.\u201d Harvard Divinity School, 2025.",
            "Ajose, Toyin Samuel. \u201cLiturgical traditions and transitions: Congregational "
            "hymn singing in Nigerian Pentecostal churches.\u201d Journal of the Association "
            "of Nigerian Musicologists, 2025.",
            "Ajose, Toyin Samuel. \u201c\u2018Me I no go suffer\u2019: Christian songs and "
            "prosperity gospel among Yoruba Pentecostals in Southwest Nigeria.\u201d",
            "Ayodeji, Oluwafemi Emmanuel. Ecstasy, Holiness and Spiritual Warfare: "
            "Yor\u00f9b\u00e1 Pentecostal Music Experience. Durham University, 2019.",
            "Course materials and lecture framing, African Urban Worship Music, Yale "
            "Divinity School.",
            "King, Roberta Rose, Jean Ngoya Kidula, James R. Krabill, and Thomas Oduro. "
            "Music in the Life of the African Church. Waco, TX: Baylor University Press, 2008.",
            "Rotimi, Oti Alaba. \u201cThe Role of Music in African Pentecostal Churches in "
            "Southwestern Nigeria.\u201d",
            "Sinach. \u201cI Know Who I Am.\u201d Official music video, 2014. YouTube.",
        ],
        "notes": (
            "Leave this slide up during questions. The core scholarly anchors are King et al. "
            "for the course framing and the Ajose and Ayodeji studies for Nigerian Pentecostal "
            "music specifically. Rotimi and Ajose\u2019s prosperity gospel essay both connect "
            "directly to the practical-theology argument."
        ),
    },
]

SCRIPT = [
    (
        "Slide 1",
        [
            "Good afternoon. The song I want to spend these eight minutes with is "
            "\u201cI Know Who I Am\u201d by Sinach. My argument is this. It is not just a song "
            "about identity. It is a practical theological act. By teaching worshippers who they "
            "are before God, the song gives them confidence, dignity, and a sense of "
            "responsibility that can move belief into action. That move from belief to action is "
            "what I want to show you."
        ],
    ),
    (
        "Slide 2",
        [
            "When you hear \u201cI Know Who I Am\u201d in a church in Lagos, Nairobi, Accra, "
            "Houston, or London, people do not just listen. They sing. They raise their hands. "
            "They repeat the same lines until those lines feel like their own.",
            "That repetition is practical theology happening in real time. It is not theology as "
            "theory. It is theology as practice. The song is training people, in the act of "
            "worship, to carry a new identity out the door with them on Monday morning. That is "
            "the thread I want to follow through the rest of the talk.",
        ],
    ),
    (
        "Slide 3",
        [
            "Quick context. Sinach is a Nigerian gospel artist associated with Christ Embassy, "
            "one of the large Pentecostal networks on the continent. \u201cI Know Who I Am\u201d "
            "came out in 2014 and has moved well beyond Nigeria. You hear it in Kenyan worship "
            "services, in Ghanaian Sunday gatherings, in African diaspora churches in the United "
            "States and Europe, and in non African churches that simply picked it up because it "
            "works.",
            "The song is a good example of African worship music traveling transnationally. And "
            "critically for my argument, it was built for participation. It was never designed "
            "to sit on a shelf. It was designed to be sung by a whole room together. That matters "
            "because formation is what a room full of voices does to you. So let me start with "
            "what the song is doing, musically.",
        ],
    ),
    (
        "Slide 4",
        [
            "Four features stand out. First, repetition. The chorus returns to the same claims "
            "over and over. Second, call and response. A worship leader sets the line, the "
            "congregation answers, the energy builds. Third, rhythm and groove. This is not a "
            "slow hymn. It moves. People clap, sway, and step with it. Fourth, the worship is "
            "embodied. You are not singing it from your head. You are singing it with your whole "
            "body.",
            "These features are not decoration. They are the machinery that makes this song "
            "formative rather than merely expressive. Now, with that machinery in mind, I want "
            "to look at what the song is actually claiming.",
        ],
    ),
    (
        "Slide 5",
        [
            "Here is what the song says. The lyrics are plain and first person. \u201cI know who "
            "God says I am.\u201d \u201cI am holy.\u201d \u201cI am righteous.\u201d \u201cI am "
            "so rich.\u201d These are not descriptions of someone else. They are claims a "
            "worshipper makes about themselves, out loud, and the room makes them together.",
            "That is already unusual. Most of the identities we carry around with us are handed "
            "to us by other people. In this song, worshippers speak a new identity in their own "
            "voice, in public, and hear everyone else speak it back. The identity being claimed "
            "is communal, embodied, and on its way to becoming practice. And that is where I "
            "want to push the argument further.",
        ],
    ),
    (
        "Slide 6",
        [
            "This is the center of the talk. What we are watching in \u201cI Know Who I Am\u201d "
            "is practical theology. Let me break that into three moves.",
            "First, the song names identity. It gives worshippers specific language about who "
            "they are before God: holy, righteous, loved, chosen. That is not abstract theology. "
            "That is theology handed to people in a form they can actually use.",
            "Second, the song forms disposition. Singing those words week after week shapes how "
            "people carry themselves. Confidence, dignity, and a sense of what they owe each "
            "other stop being moods and become habits. The congregation is rehearsing a way of "
            "standing in the world.",
            "Third, the song shapes practice. Identity that is sung becomes identity that is "
            "lived. Worshippers leave the room and show up at work, at home, and in public life "
            "with a different story about who they are. That is worship as formation. Belief "
            "that becomes action.",
        ],
    ),
    (
        "Slide 7",
        [
            "So why does this matter beyond the song? Three reasons. First, it changes how we "
            "talk about music in Christian life. Music is not the warm up before the sermon. "
            "It can be the teaching. It can shape belief, not just decorate it.",
            "Second, this is belief in action. If I have stood in a room full of people and "
            "said I am holy, I am righteous, I am loved, I have a harder time walking out and "
            "acting as if none of that is true. The song gives worshippers moral agency. It "
            "raises the floor for how they show up the rest of the week.",
            "Third, it is public Christian witness. A congregation that sings this song together "
            "does not only encourage each other on Sunday. They move into Monday with a shared "
            "confidence, a shared sense of dignity, and a shared responsibility to one another "
            "and to the world around them. One song, sung together, becomes a practiced way of "
            "being Christian in public.",
        ],
    ),
    (
        "Slide 8",
        [
            "So to close. \u201cI Know Who I Am\u201d is a worship song, yes. But it is more "
            "than that. It is a piece of practical theology set to a beat. Every time a "
            "congregation sings it, they are being formed. They are learning to carry a new "
            "identity and a new responsibility out into their lives.",
            "That is what African urban worship music does at its best. It does not only "
            "express faith. It trains people to live it, and it sends them out to live it "
            "together. I think that is why one song, written by one Nigerian artist, has ended "
            "up shaping Sunday mornings, and then Monday mornings, all over the world.",
        ],
    ),
    (
        "Slide 9",
        [
            "Thank you. The sources I used are on the slide, and I am happy to take questions.",
        ],
    ),
]

BIBLIOGRAPHY = [
    "Ajose, Toyin. \u201cMultitracking the Spirit: Musicality, Prayer Chants and Pentecostal "
    "Spirituality in Nigeria.\u201d Harvard Divinity School, 2025.",
    "Ajose, Toyin Samuel. \u201cLiturgical traditions and transitions: Congregational hymn "
    "singing in Nigerian Pentecostal churches.\u201d Journal of the Association of Nigerian "
    "Musicologists, 2025.",
    "Ajose, Toyin Samuel. \u201c\u2018Me I no go suffer\u2019: Christian songs and prosperity "
    "gospel among Yoruba Pentecostals in Southwest Nigeria.\u201d",
    "Ayodeji, Oluwafemi Emmanuel. Ecstasy, Holiness and Spiritual Warfare: Yor\u00f9b\u00e1 "
    "Pentecostal Music Experience. Durham University, 2019.",
    "Course materials and lecture framing. African Urban Worship Music. Yale Divinity School.",
    "King, Roberta Rose, Jean Ngoya Kidula, James R. Krabill, and Thomas Oduro. Music in the "
    "Life of the African Church. Waco, TX: Baylor University Press, 2008.",
    "Rotimi, Oti Alaba. \u201cThe Role of Music in African Pentecostal Churches in Southwestern "
    "Nigeria.\u201d",
    "Sinach. \u201cI Know Who I Am.\u201d Official music video, 2014. "
    "https://www.youtube.com/watch?v=frtZ4XfoXxM.",
]

# ---------------------------------------------------------------------------
# PowerPoint palette and helpers
# ---------------------------------------------------------------------------

BG        = RGBColor(0xFA, 0xF7, 0xF2)
BG_ALT    = RGBColor(0xF1, 0xEC, 0xE2)
INK       = RGBColor(0x1B, 0x1A, 0x18)
INK_SOFT  = RGBColor(0x3A, 0x36, 0x30)
MUTED     = RGBColor(0x6B, 0x63, 0x58)
ACCENT    = RGBColor(0x8A, 0x4B, 0x2A)
GOLD      = RGBColor(0xC6, 0x9B, 0x52)
PT_BG     = RGBColor(0x2A, 0x1E, 0x17)   # deep espresso for practical-theology slide
PT_CARD   = RGBColor(0x3A, 0x2A, 0x1F)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SERIF = "Georgia"
SANS  = "Calibri"


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *,
             font=SANS, size=18, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    lines = [text] if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_eyebrow(slide, left, top, text, *, color=ACCENT, width=Inches(11)):
    return add_text(
        slide, left, top, width, Inches(0.35),
        text.upper(),
        font=SANS, size=12, bold=True, color=color, line_spacing=1.0,
    )


def add_accent_bar(slide, left, top, width=Inches(0.6), height=Inches(0.05), color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_footer(slide, page_num, total, *, dark=False):
    fg = MUTED if not dark else RGBColor(0xBF, 0xB8, 0xAB)
    add_text(
        slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
        f"{TITLE}  \u00b7  {PRESENTER}",
        font=SANS, size=10, color=fg, line_spacing=1.0,
    )
    add_text(
        slide, Inches(11.7), Inches(7.05), Inches(1.0), Inches(0.3),
        f"{page_num:02d} / {total:02d}",
        font=SANS, size=10, color=fg, align=PP_ALIGN.RIGHT, line_spacing=1.0,
    )


def set_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = SANS
    run.font.size = Pt(12)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_title_slide(prs, data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)

    accent_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(0), Inches(3.833), Inches(7.5)
    )
    accent_block.fill.solid()
    accent_block.fill.fore_color.rgb = BG_ALT
    accent_block.line.fill.background()

    ring = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(10.6), Inches(0.9), Inches(1.6), Inches(1.6)
    )
    ring.fill.background()
    ring.line.color.rgb = GOLD
    ring.line.width = Pt(2)

    inner = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(10.9), Inches(1.2), Inches(1.0), Inches(1.0)
    )
    inner.fill.solid()
    inner.fill.fore_color.rgb = ACCENT
    inner.line.fill.background()

    add_eyebrow(slide, Inches(0.8), Inches(0.9), data["eyebrow"])
    add_text(
        slide, Inches(0.8), Inches(1.4), Inches(9), Inches(2.6),
        data["title"],
        font=SERIF, size=66, bold=True, color=INK, line_spacing=1.0,
    )
    add_text(
        slide, Inches(0.8), Inches(3.8), Inches(8.5), Inches(1.4),
        data["subtitle"],
        font=SERIF, size=22, italic=True, color=INK_SOFT, line_spacing=1.25,
    )

    meta_top = Inches(5.4)
    col_w = Inches(2.1)
    for i, (label, value) in enumerate(data["meta"]):
        left = Inches(0.8 + i * 2.15)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, meta_top, col_w, Inches(0.02)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = INK
        bar.line.fill.background()
        add_text(
            slide, left, meta_top + Inches(0.1), col_w, Inches(0.3),
            label.upper(),
            font=SANS, size=10, bold=True, color=MUTED, line_spacing=1.0,
        )
        add_text(
            slide, left, meta_top + Inches(0.42), col_w, Inches(0.5),
            value,
            font=SERIF, size=14, color=INK, line_spacing=1.1,
        )

    add_footer(slide, page, total)
    set_notes(slide, data["notes"])


def build_content_slide(prs, data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)

    add_eyebrow(slide, Inches(0.8), Inches(0.6), data["eyebrow"])
    add_accent_bar(slide, Inches(0.8), Inches(0.95))
    add_text(
        slide, Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.4),
        data["title"],
        font=SERIF, size=40, bold=True, color=INK, line_spacing=1.05,
    )

    top = Inches(2.9)
    gap = Inches(0.95)
    for i, bullet in enumerate(data["bullets"]):
        y = top + gap * i
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.22), Inches(0.12), Inches(0.12)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_text(
            slide, Inches(1.25), y, Inches(11.3), Inches(0.9),
            bullet,
            font=SERIF, size=22, color=INK_SOFT, line_spacing=1.25,
        )

    add_footer(slide, page, total)
    set_notes(slide, data["notes"])


def build_lyrics_slide(prs, data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)

    add_eyebrow(slide, Inches(0.8), Inches(0.6), data["eyebrow"])
    add_accent_bar(slide, Inches(0.8), Inches(0.95))
    add_text(
        slide, Inches(0.8), Inches(1.15), Inches(11.7), Inches(1.4),
        data["title"],
        font=SERIF, size=36, bold=True, color=INK, line_spacing=1.05,
    )

    panel_top = Inches(3.0)
    panel_left = Inches(0.8)
    panel_w = Inches(6.2)
    panel_h = Inches(3.4)

    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, panel_left, panel_top, panel_w, panel_h
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = BG_ALT
    panel.line.fill.background()

    gold_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, panel_left, panel_top, Inches(0.08), panel_h
    )
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = GOLD
    gold_bar.line.fill.background()

    lyrics_box = slide.shapes.add_textbox(
        panel_left + Inches(0.4), panel_top + Inches(0.3),
        panel_w - Inches(0.6), panel_h - Inches(0.6),
    )
    tf = lyrics_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(data["lyrics"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.5
        run = p.add_run()
        run.text = line
        run.font.name = SERIF
        run.font.size = Pt(22)
        run.font.italic = True
        run.font.color.rgb = INK

    add_text(
        slide, Inches(7.4), panel_top + Inches(0.1), Inches(5.2), Inches(3.4),
        data["argument"],
        font=SERIF, size=20, color=INK_SOFT, line_spacing=1.35,
    )

    add_footer(slide, page, total)
    set_notes(slide, data["notes"])


def build_practical_theology_slide(prs, data, page, total):
    """Visually distinct slide: deep background, centered title, three numbered cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PT_BG)

    # gold horizontal accent above the title
    add_accent_bar(
        slide, Inches(6.37), Inches(0.75), width=Inches(0.6), height=Inches(0.05),
        color=GOLD,
    )

    add_text(
        slide, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.4),
        data["eyebrow"].upper(),
        font=SANS, size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
        line_spacing=1.0,
    )
    add_text(
        slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.0),
        data["title"],
        font=SERIF, size=44, bold=True, color=BG, align=PP_ALIGN.CENTER,
        line_spacing=1.05,
    )
    add_text(
        slide, Inches(2.0), Inches(2.45), Inches(9.3), Inches(1.0),
        data["subtitle"],
        font=SERIF, size=18, italic=True, color=RGBColor(0xD9, 0xCE, 0xBD),
        align=PP_ALIGN.CENTER, line_spacing=1.35,
    )

    # three cards
    card_top = Inches(3.9)
    card_h = Inches(2.7)
    margin = 0.6
    gap = 0.25
    usable_w = 13.333 - (2 * margin)
    card_w = (usable_w - 2 * gap) / 3

    for i, step in enumerate(data["steps"]):
        left = Inches(margin + i * (card_w + gap))
        w = Inches(card_w)

        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, card_top, w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = PT_CARD
        card.line.fill.background()

        # gold top edge
        top_edge = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, card_top, w, Inches(0.06)
        )
        top_edge.fill.solid()
        top_edge.fill.fore_color.rgb = GOLD
        top_edge.line.fill.background()

        # big numeral
        add_text(
            slide, left + Inches(0.4), card_top + Inches(0.25),
            w - Inches(0.8), Inches(0.9),
            step["num"],
            font=SERIF, size=46, bold=True, color=GOLD, line_spacing=1.0,
        )
        # label
        add_text(
            slide, left + Inches(0.4), card_top + Inches(1.1),
            w - Inches(0.8), Inches(0.5),
            step["label"],
            font=SANS, size=16, bold=True,
            color=BG, line_spacing=1.1,
        )
        # body
        add_text(
            slide, left + Inches(0.4), card_top + Inches(1.65),
            w - Inches(0.8), card_h - Inches(1.8),
            step["body"],
            font=SERIF, size=14, italic=False,
            color=RGBColor(0xD9, 0xCE, 0xBD), line_spacing=1.35,
        )

    add_footer(slide, page, total, dark=True)
    set_notes(slide, data["notes"])


def build_quote_slide(prs, data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, INK)

    add_eyebrow(slide, Inches(0.8), Inches(0.8), data["eyebrow"], color=GOLD)
    add_text(
        slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.0),
        data["title"],
        font=SERIF, size=36, bold=True, color=BG, line_spacing=1.1,
    )

    add_text(
        slide, Inches(1.2), Inches(2.8), Inches(11.0), Inches(3.5),
        "\u201c" + data["quote"] + "\u201d",
        font=SERIF, size=28, italic=True, color=BG, line_spacing=1.35,
    )

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(6.3), Inches(0.6), Inches(0.04)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    add_text(
        slide, Inches(1.95), Inches(6.15), Inches(8), Inches(0.4),
        f"Thesis close  \u00b7  {PRESENTER}",
        font=SANS, size=12, color=RGBColor(0xBF, 0xB8, 0xAB), line_spacing=1.0,
    )

    add_footer(slide, page, total, dark=True)
    set_notes(slide, data["notes"])


def build_bibliography_slide(prs, data, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)

    add_eyebrow(slide, Inches(0.8), Inches(0.6), data["eyebrow"])
    add_accent_bar(slide, Inches(0.8), Inches(0.95))
    add_text(
        slide, Inches(0.8), Inches(1.15), Inches(11), Inches(1.2),
        data["title"],
        font=SERIF, size=40, bold=True, color=INK, line_spacing=1.05,
    )

    # Auto-scale type size so long bibliographies still fit on one slide.
    n = len(data["entries"])
    if n <= 4:
        num_pt, body_pt, line_sp, space_after_pt = 14, 15, 1.3, 8
    elif n <= 6:
        num_pt, body_pt, line_sp, space_after_pt = 12, 13, 1.25, 6
    else:
        num_pt, body_pt, line_sp, space_after_pt = 11, 12, 1.2, 4

    box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.55), Inches(11.7), Inches(4.35)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, entry in enumerate(data["entries"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_sp
        p.space_after = Pt(space_after_pt)

        number_run = p.add_run()
        number_run.text = f"{i + 1}.  "
        number_run.font.name = SANS
        number_run.font.size = Pt(num_pt)
        number_run.font.bold = True
        number_run.font.color.rgb = ACCENT

        text_run = p.add_run()
        text_run.text = entry
        text_run.font.name = SERIF
        text_run.font.size = Pt(body_pt)
        text_run.font.color.rgb = INK_SOFT

    add_footer(slide, page, total)
    set_notes(slide, data["notes"])


BUILDERS = {
    "title": build_title_slide,
    "content": build_content_slide,
    "lyrics": build_lyrics_slide,
    "practical_theology": build_practical_theology_slide,
    "quote": build_quote_slide,
    "bibliography": build_bibliography_slide,
}


def build_pptx(path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = len(SLIDES)
    for i, data in enumerate(SLIDES, start=1):
        BUILDERS[data["kind"]](prs, data, i, total)

    prs.save(path)


# ---------------------------------------------------------------------------
# Word script
# ---------------------------------------------------------------------------

def _set_font(run, name=None, size=None, bold=None, italic=None, color=None):
    if name is not None:
        run.font.name = name
        rPr = run._element.get_or_add_rPr()
        rFonts = rPr.find(qn("w:rFonts"))
        if rFonts is None:
            rFonts = etree.SubElement(rPr, qn("w:rFonts"))
        rFonts.set(qn("w:ascii"), name)
        rFonts.set(qn("w:hAnsi"), name)
        rFonts.set(qn("w:cs"), name)
    if size is not None:
        run.font.size = DocxPt(size)
    if bold is not None:
        run.bold = bold
    if italic is not None:
        run.italic = italic
    if color is not None:
        run.font.color.rgb = color


def build_docx(path: str):
    doc = Document()

    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = DocxPt(12)

    for section in doc.sections:
        section.top_margin = DocxInches(1.0)
        section.bottom_margin = DocxInches(1.0)
        section.left_margin = DocxInches(1.1)
        section.right_margin = DocxInches(1.1)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r = title.add_run(TITLE)
    _set_font(r, name="Georgia", size=28, bold=True, color=DocxRGB(0x1B, 0x1A, 0x18))

    sub = doc.add_paragraph()
    r = sub.add_run(SUBTITLE)
    _set_font(r, name="Georgia", size=14, italic=True, color=DocxRGB(0x3A, 0x36, 0x30))

    meta = doc.add_paragraph()
    meta_text = f"Presenter: {PRESENTER}    |    Course: {COURSE}    |    Target length: 8 minutes"
    r = meta.add_run(meta_text)
    _set_font(r, name="Calibri", size=11, color=DocxRGB(0x6B, 0x63, 0x58))

    h = doc.add_paragraph()
    r = h.add_run("Thesis")
    _set_font(r, name="Georgia", size=14, bold=True, color=DocxRGB(0x8A, 0x4B, 0x2A))

    thesis = doc.add_paragraph()
    r = thesis.add_run(THESIS)
    _set_font(r, name="Georgia", size=12, italic=True, color=DocxRGB(0x1B, 0x1A, 0x18))

    doc.add_paragraph()

    h = doc.add_paragraph()
    r = h.add_run("Full script")
    _set_font(r, name="Georgia", size=14, bold=True, color=DocxRGB(0x8A, 0x4B, 0x2A))

    note = doc.add_paragraph()
    r = note.add_run(
        "The script below runs about 8 minutes at a natural pace. Slide cues are bracketed so "
        "you can advance without losing your place. The argument moves from identity to action, "
        "with the practical theology slide as the hinge."
    )
    _set_font(r, name="Calibri", size=11, italic=True, color=DocxRGB(0x6B, 0x63, 0x58))

    word_count = 0
    for cue, paragraphs in SCRIPT:
        cue_p = doc.add_paragraph()
        cue_p.paragraph_format.space_before = DocxPt(12)
        r = cue_p.add_run(f"[{cue}]")
        _set_font(r, name="Calibri", size=11, bold=True, color=DocxRGB(0x8A, 0x4B, 0x2A))

        for para in paragraphs:
            p = doc.add_paragraph()
            p.paragraph_format.space_after = DocxPt(6)
            r = p.add_run(para)
            _set_font(r, name="Calibri", size=12, color=DocxRGB(0x1B, 0x1A, 0x18))
            word_count += len(para.split())

    doc.add_paragraph()
    wc = doc.add_paragraph()
    r = wc.add_run(f"Approximate spoken word count: {word_count} words.")
    _set_font(r, name="Calibri", size=10, italic=True, color=DocxRGB(0x6B, 0x63, 0x58))

    doc.add_page_break()
    h = doc.add_paragraph()
    r = h.add_run("Bibliography")
    _set_font(r, name="Georgia", size=16, bold=True, color=DocxRGB(0x1B, 0x1A, 0x18))

    for entry in BIBLIOGRAPHY:
        p = doc.add_paragraph(style="List Number")
        r = p.add_run(entry)
        _set_font(r, name="Calibri", size=11, color=DocxRGB(0x1B, 0x1A, 0x18))
        if entry.lower().startswith("[verify"):
            r.italic = True
            r.font.color.rgb = DocxRGB(0xA8, 0x5A, 0x2C)

    doc.save(path)
    return word_count


# ---------------------------------------------------------------------------
# Entrypoint
# ---------------------------------------------------------------------------

def _safe_write(primary, writer):
    """Write to primary path. If locked (file open), fall back to .new.<ext>."""
    try:
        writer(primary)
        return primary
    except PermissionError:
        stem, _, ext = primary.rpartition(".")
        fallback = f"{stem}.new.{ext}"
        writer(fallback)
        return fallback


def main():
    pptx_path = "I_Know_Who_I_Am_Presentation.pptx"
    docx_path = "I_Know_Who_I_Am_Full_Script.docx"

    pptx_out = _safe_write(pptx_path, build_pptx)

    wc_box = {}
    def _write_docx(path):
        wc_box["wc"] = build_docx(path)
    docx_out = _safe_write(docx_path, _write_docx)

    print(f"Wrote {pptx_out}")
    print(f"Wrote {docx_out}  (script word count: {wc_box['wc']})")
    if pptx_out != pptx_path or docx_out != docx_path:
        print(
            "\nNote: one or more original files were locked (probably open in "
            "PowerPoint/Word). Close them and re-run to overwrite, or rename the "
            "'.new.' files manually."
        )


if __name__ == "__main__":
    main()
